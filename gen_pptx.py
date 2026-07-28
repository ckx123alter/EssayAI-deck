"""
将 HTML 幻灯片 Deck 转换为 PPTX 格式

用法：python gen_pptx.py
输出：EssayAI_产品演示.pptx
"""
import os
import sys
import tempfile
import time
from pathlib import Path

# 修复 Windows GBK 终端下 emoji 输出问题
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

# ── 配置 ──────────────────────────────────────────────
BASE_DIR = Path(os.path.dirname(os.path.abspath(__file__)))
SLIDES_DIR = BASE_DIR / "slides"
VIDEO_PATH = BASE_DIR / "演示视频.mp4"
OUTPUT_PATH = BASE_DIR / "EssayAI_产品演示.pptx"

SLIDE_W = 1920   # 设计稿宽度
SLIDE_H = 1080   # 设计稿高度
SCALE = 2        # 截图缩放因子（2x 视网膜清晰度）
ANIM_WAIT = 5    # 等待动画完成的秒数

# 幻灯片顺序（与 index.html MANIFEST 一致）
MANIFEST = [
    ("01-cover.html",      "封面 · EssayAI"),
    ("02-investment.html", "项目概要"),
    ("03-market.html",     "市场洞察"),
    ("04-painpoints.html", "用户痛点"),
    ("05-solution.html",   "产品方案"),
    ("06-tech.html",       "核心技术"),
    ("07-demo.html",       "产品效果"),
    ("08-moat.html",       "竞品对比"),
    ("09-roadmap.html",    "增长路线"),
    ("10-financials.html", "财务预测"),
    ("11-funding.html",    "融资需求"),
    ("12-team.html",       "核心团队"),
    ("13-thanks.html",     "谢谢"),
]

# ── 工具函数 ──────────────────────────────────────────
def screenshot_slides(playwright) -> list[Path]:
    """用 Chromium 逐一渲染并截图每页幻灯片，返回截图路径列表。"""
    browser = playwright.chromium.launch()
    context = browser.new_context(
        viewport={"width": SLIDE_W, "height": SLIDE_H},
        device_scale_factor=SCALE,
    )
    page = context.new_page()

    tmp_dir = Path(tempfile.mkdtemp(prefix="essayai_slides_"))
    screenshots = []

    for idx, (filename, label) in enumerate(MANIFEST):
        html_path = SLIDES_DIR / filename
        file_url = html_path.as_uri()

        print(f"  [{idx+1:2d}/13] 截图中: {label} ...", end=" ", flush=True)

        page.goto(file_url, wait_until="networkidle", timeout=30000)

        # 等待所有 CSS 动画播放完毕 + 字体加载
        try:
            page.wait_for_function(
                "() => document.fonts ? document.fonts.ready.then(() => true) : true",
                timeout=8000,
            )
        except Exception:
            pass  # 离线时字体可能加载失败，不阻塞

        # 额外等待确保 animation-delay 最晚的元素已显示
        # 各页动画最晚 ~2.6s（07-demo.html），给 5s 余量
        page.wait_for_timeout(ANIM_WAIT * 1000)

        out_path = tmp_dir / f"slide_{idx+1:02d}.png"
        page.screenshot(path=str(out_path), full_page=False)
        screenshots.append(out_path)

        print("✓")

    browser.close()
    return screenshots


def create_pptx(screenshots: list[Path]):
    """基于截图创建 PPTX 文件。"""
    prs = Presentation()
    # 设为 16:9 宽屏（与 1920×1080 比例一致）
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout

    for idx, (screenshot_path, (filename, label)) in enumerate(
        zip(screenshots, MANIFEST)
    ):
        slide = prs.slides.add_slide(blank_layout)

        # ── 全屏背景截图 ──
        slide.shapes.add_picture(
            str(screenshot_path),
            left=Inches(0),
            top=Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # ── 第 7 页：嵌入演示视频 ──
        if filename == "07-demo.html" and VIDEO_PATH.exists():
            # 视频放在右下角小区域，覆盖在截图上
            video_left = Inches(8.6)
            video_top = Inches(5.6)
            video_width = Inches(4.4)
            video_height = Inches(1.7)

            try:
                movie = slide.shapes.add_movie(
                    str(VIDEO_PATH),
                    left=video_left,
                    top=video_top,
                    width=video_width,
                    height=video_height,
                    poster_frame_image=None,
                )
                print(f"  ✓ 视频已嵌入第 {idx+1} 页")
            except Exception as e:
                print(f"  ⚠ 视频嵌入失败 (第 {idx+1} 页): {e}")

        # ── 页码标签（右下角） ──
        page_num = idx + 1
        txBox = slide.shapes.add_textbox(
            left=Inches(12.2),
            top=Inches(7.0),
            width=Inches(0.9),
            height=Inches(0.35),
        )
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"{page_num} / 13"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p.font.name = "Inter"
        p.alignment = PP_ALIGN.RIGHT

        print(f"  [{idx+1:2d}/13] ✓ 已添加: {label}")

    # ── 保存 ──
    prs.save(str(OUTPUT_PATH))
    print(f"\n✅ 完成: {OUTPUT_PATH}")


# ── 主流程 ────────────────────────────────────────────
def main():
    print("=" * 60)
    print("EssayAI HTML Deck → PPTX 转换器")
    print("=" * 60)
    print()

    # 检查 slides 目录
    if not SLIDES_DIR.exists():
        print(f"❌ 错误: slides 目录不存在: {SLIDES_DIR}")
        sys.exit(1)

    if not VIDEO_PATH.exists():
        print(f"⚠ 警告: 演示视频不存在，将跳过视频嵌入: {VIDEO_PATH}")

    t0 = time.time()

    # Phase 1: 截图
    print("📷 Phase 1/2: 渲染截图 (Chromium headless)")
    with sync_playwright() as p:
        screenshots = screenshot_slides(p)

    if not screenshots:
        print("❌ 错误: 截图失败，未生成任何图片")
        sys.exit(1)

    print(f"\n  共 {len(screenshots)} 张截图，耗时 {time.time()-t0:.1f}s\n")

    # Phase 2: 生成 PPTX
    t1 = time.time()
    print("📦 Phase 2/2: 生成 PPTX")
    create_pptx(screenshots)
    print(f"  耗时 {time.time()-t1:.1f}s")

    # 清理临时截图
    tmp_dir = screenshots[0].parent
    for s in screenshots:
        s.unlink(missing_ok=True)
    tmp_dir.rmdir()
    print(f"\n🧹 已清理临时文件")

    print(f"\n总耗时 {time.time()-t0:.1f}s")


if __name__ == "__main__":
    main()
